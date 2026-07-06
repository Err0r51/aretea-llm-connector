"""AC4: extract_pptx returns per-slide text with slide headers (real .pptx)."""

from __future__ import annotations

import io

from aretea_drive_mcp.gdrive.extract import extract_pptx
from pptx import Presentation

CEIL = 314_572_800


def _pptx_bytes() -> bytes:
    prs = Presentation()
    for title in ("ALPHA SLIDE", "BETA SLIDE"):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # "Title Only" layout
        slide.shapes.title.text = title
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_returns_per_slide_text() -> None:
    text = extract_pptx(_pptx_bytes(), max_uncompressed_bytes=CEIL)
    assert "--- Slide 1 ---" in text
    assert "--- Slide 2 ---" in text
    assert "ALPHA SLIDE" in text
    assert "BETA SLIDE" in text
