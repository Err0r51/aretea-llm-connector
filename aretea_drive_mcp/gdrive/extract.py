"""Pure in-process text extraction for binary Drive files (PRD 5).

Text-layer **PDF** and modern OOXML **.docx / .xlsx / .pptx** are parsed from their raw bytes — the
bytes are downloaded read-only via ``get_media`` in ``client.py`` and handed here. There is **no
Google-side conversion**: converting an Office file to a Google Doc to ``export`` it would be a
mutation (``files.copy``) and need an SA write scope — both forbidden by the connector's read-only /
no-delegation invariants (PRD 2 AC8/AC9). Local parsing is the only path consistent with that model.

Every function is a pure ``bytes -> str``: no network, no event loop, no ``DriveClient`` — which is
exactly what makes them unit-testable against real fixture files without mocking Google. They return
extracted text on success and raise **only** :class:`ExtractionError` for the expected-degraded
cases (scanned/no-text-layer PDF, encrypted, corrupt, decompression bomb). The caller turns that
into a *stated* placeholder (never-silent, mirroring PRD 2 AC11); it never logs here (logging lives
in the dispatcher, to keep these functions pure).

⚠ Library call shapes (``PdfReader``, python-docx body walk, ``load_workbook(data_only=...)``,
python-pptx shape walk) are written to the pinned versions; verify at build if the pins move.

.doc/.xls/.ppt (pre-2007 OLE) and OCR for scanned PDFs are out of scope (Non-Goals) — handled as
stated placeholders by the caller, not here.
"""

from __future__ import annotations

import zipfile
from io import BytesIO

from docx import Document
from docx.oxml.ns import qn
from docx.table import Table
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

# --- stated placeholders for expected-degraded content (never-silent; PRD 2 AC11) ---
NO_TEXT_LAYER_PLACEHOLDER = (
    "[PDF has no extractable text layer — likely scanned; OCR not supported in v1]"
)
ENCRYPTED_PLACEHOLDER = "[file is encrypted / password-protected — cannot extract text]"
CORRUPT_PLACEHOLDER = "[file appears corrupt or is not the expected format — could not parse]"
BOMB_PLACEHOLDER = (
    "[file's uncompressed size exceeds the extraction limit — refused (possible zip bomb)]"
)


class ExtractionError(Exception):
    """A narrow, *expected* extraction failure — not a bug.

    ``reason`` is the machine token the dispatcher logs as ``outcome`` and is one of
    ``{"no_text_layer", "encrypted", "corrupt", "decompression_bomb_refused"}``. ``placeholder`` is
    the stated string returned to the caller in place of content.
    """

    def __init__(self, reason: str, placeholder: str) -> None:
        super().__init__(reason)
        self.reason = reason
        self.placeholder = placeholder


def _guard_ooxml_bomb(data: bytes, ceiling: int) -> None:
    """Refuse an OOXML zip whose declared **uncompressed** size exceeds ``ceiling``.

    OOXML files are ZIP containers; a small archive can inflate to gigabytes. We sum the
    central-directory ``file_size`` (uncompressed, reported without decompressing) as a cheap
    pre-check. This is paired with the caller's compressed-input cap (``max_input_bytes``), which
    bounds what can even be delivered; a crafted archive that *lies* about ``file_size`` is a
    residual, accepted risk beyond this guard. A non-zip payload here means a corrupt file.
    """
    try:
        with zipfile.ZipFile(BytesIO(data)) as zf:
            total_uncompressed = sum(info.file_size for info in zf.infolist())
    except zipfile.BadZipFile as err:
        raise ExtractionError("corrupt", CORRUPT_PLACEHOLDER) from err
    if total_uncompressed > ceiling:
        raise ExtractionError("decompression_bomb_refused", BOMB_PLACEHOLDER)


def extract_pdf(data: bytes, *, max_uncompressed_bytes: int) -> str:
    """Text-layer PDF → concatenated page text. ``max_uncompressed_bytes`` unused (uniform sig)."""
    try:
        reader = PdfReader(BytesIO(data))
    except Exception as err:  # malformed / empty / not-a-PDF → degrade, don't crash
        raise ExtractionError("corrupt", CORRUPT_PLACEHOLDER) from err
    if reader.is_encrypted:
        raise ExtractionError("encrypted", ENCRYPTED_PLACEHOLDER)
    parts: list[str] = []
    try:
        for page in reader.pages:
            parts.append(str(page.extract_text() or ""))
    except Exception as err:  # a page-level parse failure means the PDF is malformed
        raise ExtractionError("corrupt", CORRUPT_PLACEHOLDER) from err
    text = "\n".join(parts)
    if not text.strip():
        # A born-scanned PDF has pages but no text layer — say so; OCR is out of scope.
        raise ExtractionError("no_text_layer", NO_TEXT_LAYER_PLACEHOLDER)
    return text


def extract_docx(data: bytes, *, max_uncompressed_bytes: int) -> str:
    """.docx → paragraph and table-cell text, in **document order** (body walk).

    ``document.paragraphs`` and ``document.tables`` are separate collections that lose interleaving,
    so we walk the body's child elements directly. Headers/footers/textboxes/comments are out of
    scope (basic fidelity).
    """
    _guard_ooxml_bomb(data, max_uncompressed_bytes)
    try:
        doc = Document(BytesIO(data))
    except Exception as err:
        raise ExtractionError("corrupt", CORRUPT_PLACEHOLDER) from err
    lines: list[str] = []
    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            text = str(Paragraph(child, doc).text)
            if text.strip():
                lines.append(text)
        elif child.tag == qn("w:tbl"):
            for row in Table(child, doc).rows:
                lines.append("\t".join(str(cell.text) for cell in row.cells))
    return "\n".join(lines)


def extract_xlsx(data: bytes, *, max_uncompressed_bytes: int) -> str:
    """.xlsx → per-sheet cell **values** (never formula strings; mirrors the Sheets read).

    Uses ``data_only=True`` — returns each cell's last cached value. A sheet generated
    programmatically and never opened in a spreadsheet app has no cached value, so formula cells
    read back empty (documented limitation).
    """
    _guard_ooxml_bomb(data, max_uncompressed_bytes)
    try:
        wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
    except Exception as err:
        raise ExtractionError("corrupt", CORRUPT_PLACEHOLDER) from err
    try:
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f"# {ws.title}")
            for values in ws.iter_rows(values_only=True):
                parts.append("\t".join("" if cell is None else str(cell) for cell in values))
    finally:
        wb.close()
    return "\n".join(parts)


def extract_pptx(data: bytes, *, max_uncompressed_bytes: int) -> str:
    """.pptx → per-slide shape text. Grouped-shape text and speaker notes are out of scope."""
    _guard_ooxml_bomb(data, max_uncompressed_bytes)
    try:
        prs = Presentation(BytesIO(data))
    except Exception as err:
        raise ExtractionError("corrupt", CORRUPT_PLACEHOLDER) from err
    lines: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        lines.append(f"--- Slide {index} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = str(shape.text_frame.text)
                if text.strip():
                    lines.append(text.rstrip("\n"))
    return "\n".join(lines)
